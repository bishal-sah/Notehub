"""
Flashcard generation service.
Extracts text from uploaded note files (PDF, DOC/DOCX, PPT/PPTX) and
auto-generates Q&A flashcards using sentence-based keyword extraction.
"""
import os
import re
import logging
from collections import Counter

logger = logging.getLogger(__name__)


# ─── Text Extraction ─────────────────────────────────────

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file."""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return '\n'.join(pages)
    except Exception as e:
        logger.warning(f"PDF extraction failed: {e}")
        return ''


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from a DOCX file."""
    try:
        from docx import Document
        doc = Document(file_path)
        return '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        logger.warning(f"DOCX extraction failed: {e}")
        return ''


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a PPTX file (best-effort via python-pptx or fallback)."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text)
        return '\n'.join(texts)
    except ImportError:
        logger.warning("python-pptx not installed, skipping PPTX extraction")
        return ''
    except Exception as e:
        logger.warning(f"PPTX extraction failed: {e}")
        return ''


def extract_text(file_path: str, file_type: str) -> str:
    """Route to the correct extractor based on file type."""
    ft = file_type.lower()
    if ft == 'pdf':
        return extract_text_from_pdf(file_path)
    elif ft in ('doc', 'docx'):
        return extract_text_from_docx(file_path)
    elif ft in ('ppt', 'pptx'):
        return extract_text_from_pptx(file_path)
    return ''


# ─── Text Cleaning ───────────────────────────────────────

STOPWORDS = set("""
a about above after again against all am an and any are aren't as at be because
been before being below between both but by can't cannot could couldn't did didn't
do does doesn't doing don't down during each few for from further get got had hadn't
has hasn't have haven't having he he'd he'll he's her here here's hers herself him
himself his how how's i i'd i'll i'm i've if in into is isn't it it's its itself
let's me more most mustn't my myself no nor not of off on once only or other ought
our ours ourselves out over own same shan't she she'd she'll she's should shouldn't
so some such than that that's the their theirs them themselves then there there's
these they they'd they'll they're they've this those through to too under until up
very was wasn't we we'd we'll we're we've were weren't what what's when when's where
where's which while who who's whom why why's will with won't would wouldn't you you'd
you'll you're you've your yours yourself yourselves also however although therefore
thus hence moreover furthermore nevertheless nonetheless whereas etc e.g i.e eg ie
""".split())

# Extra non-concept words that appear frequently but make bad flashcard topics
JUNK_KEYWORDS = set("""
used using uses use make makes made making call called calls calling
example examples include includes including included give gives given gave
take takes taken took show shows shown showed provide provides provided
become becomes became allow allows allowed need needs needed require requires
required create creates created different various many much several just
like also well still even new first last next following way ways set
get gets getting thing things part parts case cases number numbers time
times point points result results type types work works working based
able may might often usually always never sometimes certain certain well
see seen know known find found found right left good better best also
come comes came going gone run runs running let help helps helped start
started end ends main important can one two three four five said say
another each every some any such only own same than too high low
large small big long short old new great little different used using
write writes written read reads reading try tries tried keep keeps kept
change changes changed move moves moved ask asks asked tell tells told
open opens opened close closes closed turn turns turned put puts contain
contains contained lead leads led follow follows followed begin begins began
think thinks thought look looks looked want wants wanted seem seems seemed
add adds added mean means meant look looks looked play plays played
hold holds held bring brings brought happen happens happened must shall
along within without through since before after between both either neither
nor whether while until although though really actually simply basically
currently recently already probably perhaps maybe actually simply just
quite rather enough almost nearly pretty fairly somewhat really truly
actually certainly definitely usually generally specifically particularly
especially exactly simply clearly obviously apparently overall per
""".split())


def is_garbage_text(text: str) -> bool:
    """
    Detect text that is just watermarks, repeated tokens, or scanner artefacts
    (e.g. 'CamScanner\\nCamScanner\\n...' from scanned PDFs).
    Returns True if the text should be treated as empty.
    """
    if not text or len(text.strip()) < 30:
        return True
    # Collapse to unique non-empty lines
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    if not lines:
        return True
    unique = set(lines)
    # If >80 % of lines are the same token, it's garbage
    if len(unique) <= 3 and len(lines) > 5:
        return True
    # If almost every line is a known scanner watermark
    watermarks = {'camscanner', 'cam scanner', 'scanbot', 'adobe scan',
                  'microsoft lens', 'office lens', 'genius scan', 'scanner pro'}
    wm_count = sum(1 for l in lines if l.lower() in watermarks)
    if wm_count > len(lines) * 0.5:
        return True
    return False


def clean_text(text: str) -> str:
    """Clean extracted text for processing."""
    # Remove URLs
    text = re.sub(r'https?://\S+', '', text)
    # Remove email addresses
    text = re.sub(r'\S+@\S+', '', text)
    # Remove author / metadata lines ("Notes prepared by ...", "Compiled by ...", etc.)
    text = re.sub(
        r'(?i)\b(?:notes?\s+)?(?:prepared|compiled|written|created|made|authored|submitted)'
        r'\s+by\s+[A-Z][a-zA-Z\s,\.&]+',
        '', text
    )
    # Remove standalone author attribution lines
    text = re.sub(r'(?i)^\s*(?:by|author|name)\s*[:;]?\s*[A-Z][a-zA-Z\s]+$', '', text, flags=re.MULTILINE)
    # Remove page numbers
    text = re.sub(r'(?i)\b(?:page|pg\.?)\s*\d+\b', '', text)
    text = re.sub(r'(?m)^\s*\d+\s*$', '', text)
    # Remove chapter/section numbering at start of lines (e.g. "2:", "2.1", "Chapter 3")
    text = re.sub(r'(?m)^\s*(?:chapter|unit|module|section|lecture)\s*\d+', '', text, flags=re.IGNORECASE)
    # Collapse whitespace
    text = re.sub(r'\s+', ' ', text)
    return text.strip()


def split_sentences(text: str) -> list[str]:
    """Split text into sentences."""
    # Split on sentence-ending punctuation, newline-like breaks, or dash-separated sections
    raw = re.split(r'(?<=[.!?])\s+(?=[A-Z])|\s+-\s+(?=[A-Z])', text)
    # Further split any overly long chunks on semicolons or colons followed by space+capital
    expanded: list[str] = []
    for chunk in raw:
        if len(chunk) > 300:
            sub = re.split(r'[;]\s+(?=[A-Z])', chunk)
            expanded.extend(sub)
        else:
            expanded.append(chunk)
    # Filter out very short or very long sentences
    return [
        s.strip() for s in expanded
        if 30 < len(s.strip()) < 400
    ]


def _is_junk_sentence(sentence: str) -> bool:
    """Detect bibliography, citation, reference lines, or other non-content."""
    s = sentence.strip()
    ref_patterns = [
        r'\b(?:ISBN|ISSN|DOI|pp\.|vol\.|edition|ed\.|publisher|press|published)\b',
        r'\b(?:et\s+al|ibid)\b',
        r'\d{4}\)\s*[,.]',
        r'–\s*[A-Z][a-z]+\s+[A-Z]',
        r'^\d+\.\s+[A-Z]',
        r'&\s+[A-Z][a-z]+\s+\d',
        r'^\s*(?:figure|fig\.|table|source|reference|bibliography)\s',
        r'©|copyright|all\s+rights\s+reserved',
        r'(?i)\b(?:prepared|compiled|written|authored|submitted)\s+by\b',
        r'(?i)\b(?:notes?\s+by|lecture\s+notes?)\b.*\b[A-Z][a-z]+\s+[A-Z][a-z]+',
    ]
    for pat in ref_patterns:
        if re.search(pat, s, re.IGNORECASE):
            return True
    words = s.split()
    if len(words) < 6:
        return True
    if len(words) < 8:
        caps = sum(1 for w in words if w[0:1].isupper()
                   and w not in ('The', 'A', 'An', 'In', 'Of', 'And', 'For', 'To', 'It', 'This'))
        if caps > len(words) * 0.6:
            return True
    return False


def _extract_concepts(text: str, top_n: int = 40) -> list[str]:
    """
    Extract meaningful concept keywords — nouns, technical terms, proper nouns.
    Filters out common verbs, adjectives, and generic words that make bad flashcard topics.
    Also detects multi-word terms (bigrams) like 'object oriented' or 'virtual machine'.
    """
    all_stop = STOPWORDS | JUNK_KEYWORDS

    # --- Single-word concepts ---
    words = re.findall(r'\b[a-zA-Z]{3,}\b', text)
    words_lower = [w.lower() for w in words]

    # Count frequencies of non-stopwords
    single_counts: Counter = Counter()
    for w in words_lower:
        if w not in all_stop and len(w) >= 4:
            single_counts[w] += 1

    # Boost words that appear capitalised in the text (likely proper nouns / terms)
    cap_words = re.findall(r'\b[A-Z][a-z]{2,}\b', text)
    for w in cap_words:
        wl = w.lower()
        if wl not in all_stop and wl in single_counts:
            single_counts[wl] += 1  # double-count capitalised occurrences

    # --- Multi-word concepts (bigrams) ---
    bigram_counts: Counter = Counter()
    for i in range(len(words_lower) - 1):
        a, b = words_lower[i], words_lower[i + 1]
        if a not in all_stop and b not in all_stop and len(a) >= 3 and len(b) >= 3:
            bigram = f"{a} {b}"
            bigram_counts[bigram] += 1

    # Only keep bigrams that appear at least twice
    bigrams = {bg: c for bg, c in bigram_counts.items() if c >= 2}

    # Merge: bigrams first (higher quality), then singles
    concepts: list[str] = []
    seen: set[str] = set()

    for bg, _ in sorted(bigrams.items(), key=lambda x: x[1], reverse=True):
        if len(concepts) >= top_n:
            break
        concepts.append(bg)
        seen.add(bg)
        # Mark component words as used
        for part in bg.split():
            seen.add(part)

    for word, _ in single_counts.most_common(top_n * 2):
        if len(concepts) >= top_n:
            break
        if word not in seen:
            concepts.append(word)
            seen.add(word)

    return concepts


# ─── Code / Syntax Extraction ────────────────────────────

# Common programming keywords that signal code
_CODE_KEYWORDS = re.compile(
    r'\b(?:public|private|protected|static|void|class|interface|abstract|'
    r'extends|implements|import|package|return|new|this|super|final|'
    r'int|float|double|char|boolean|string|long|byte|short|'
    r'if|else|for|while|do|switch|case|break|continue|try|catch|throw|throws|finally|'
    r'def|self|print|lambda|async|await|const|let|var|function|'
    r'System\.out|main|args|null|true|false|None|True|False)\b',
    re.IGNORECASE
)

_CODE_SYMBOLS = re.compile(r'[{};()=<>]')


def _is_code_line(line: str) -> bool:
    """Heuristic: does this line look like source code?"""
    stripped = line.strip()
    if not stripped or len(stripped) < 5:
        return False
    kw_hits = len(_CODE_KEYWORDS.findall(stripped))
    sym_hits = len(_CODE_SYMBOLS.findall(stripped))
    # A line is likely code if it has 2+ keywords or 3+ code symbols
    if kw_hits >= 2 or sym_hits >= 3:
        return True
    # Lines ending with { or ; or ) { are almost certainly code
    if re.search(r'[{;]\s*$', stripped):
        return True
    return False


def _extract_code_blocks(raw_text: str) -> list[str]:
    """
    Extract contiguous code blocks from the raw text.
    Returns a list of code snippet strings.
    """
    lines = raw_text.split('\n')
    blocks: list[str] = []
    current_block: list[str] = []

    for line in lines:
        if _is_code_line(line):
            current_block.append(line.rstrip())
        else:
            if len(current_block) >= 2:  # at least 2 lines to be a block
                blocks.append('\n'.join(current_block))
            current_block = []

    # Don't forget trailing block
    if len(current_block) >= 2:
        blocks.append('\n'.join(current_block))

    return blocks


def _find_syntax_for_concept(concept: str, code_blocks: list[str]) -> str:
    """
    Find the best code snippet that illustrates a concept.
    Returns the snippet string, or '' if none found.
    """
    concept_lower = concept.lower()
    parts = concept_lower.split()

    best_block = ''
    best_score = 0

    for block in code_blocks:
        block_lower = block.lower()
        score = 0

        # Full concept appears in the code
        if concept_lower in block_lower:
            score += 10
        # Individual words appear
        for part in parts:
            if part in block_lower:
                score += 3

        # Prefer shorter, focused snippets (not giant dumps)
        line_count = block.count('\n') + 1
        if 2 <= line_count <= 12:
            score += 5
        elif line_count <= 20:
            score += 2

        if score > best_score:
            best_score = score
            best_block = block

    # Only return if there's a meaningful match
    if best_score >= 6:
        # Trim to max ~15 lines
        lines = best_block.split('\n')
        if len(lines) > 15:
            lines = lines[:15]
            lines.append('...')
        return '\n'.join(lines)
    return ''


# ─── Flashcard Generation ────────────────────────────────

def _extract_definition(sentence: str, concept: str) -> str | None:
    """
    If the sentence defines or explains the concept, extract a clean answer.
    Returns the definition part, or None if it's not a real definition.
    """
    if _is_junk_sentence(sentence):
        return None

    s = sentence.strip()
    c_esc = re.escape(concept)

    # Pattern: "Concept is/are/was ..." → extract everything after the verb
    m = re.search(
        rf'\b{c_esc}\b\s+(?:is|are|was|were)\s+(.{{20,}})',
        s, re.IGNORECASE
    )
    if m and m.start() < 50:
        definition = m.group(0).strip()
        if not definition.endswith('.'):
            definition += '.'
        return definition

    # Pattern: "Concept refers to / means / can be defined as ..."
    m = re.search(
        rf'\b{c_esc}\b\s+(?:refers?\s+to|means?|represents?|describes?|'
        rf'can\s+be\s+defined\s+as|is\s+defined\s+as|is\s+known\s+as)\s+(.{{20,}})',
        s, re.IGNORECASE
    )
    if m and m.start() < 60:
        definition = m.group(0).strip()
        if not definition.endswith('.'):
            definition += '.'
        return definition

    # Pattern: "Concept: explanation" or "Concept – explanation" at sentence start
    m = re.match(rf'\s*{c_esc}\s*[-:–]\s+(.{{25,}})', s, re.IGNORECASE)
    if m:
        return s.strip()

    return None


def _make_concept_question(concept: str) -> str:
    """Generate a natural question for a concept based on its form."""
    c = concept.strip()
    # Multi-word or capitalised → likely a named concept
    if ' ' in c or c[0].isupper():
        return f"What is {c}?"
    return f"Define {c}."


def _make_explanation_card(sentence: str, concept: str) -> dict | None:
    """
    Create a card from a sentence that explains what a concept does / is used for,
    even if it's not a strict definition.
    E.g. "Polymorphism allows objects to take many forms."
    """
    if _is_junk_sentence(sentence):
        return None

    s = sentence.strip()
    c_esc = re.escape(concept)

    # Concept + action verb near start: "X allows/enables/provides/supports/helps ..."
    m = re.search(
        rf'\b{c_esc}\b\s+(?:allows?|enables?|provides?|supports?|helps?|'
        rf'ensures?|handles?|manages?|performs?|converts?|stores?|'
        rf'creates?|implements?|controls?|processes?)\s+(.{{15,}})',
        s, re.IGNORECASE
    )
    if m and m.start() < 50:
        answer = s
        if not answer.endswith('.'):
            answer += '.'
        if len(answer) < 40 or len(answer) > 400:
            return None
        question = f"What does {concept} do?"
        return {'question': question, 'answer': answer}

    return None


def _make_fill_blank_card(sentence: str, concept: str) -> dict | None:
    """Create a fill-in-the-blank card — only for strong conceptual sentences."""
    if _is_junk_sentence(sentence):
        return None

    s = sentence.strip()
    if len(s) < 50 or len(s) > 300:
        return None

    c_esc = re.escape(concept)
    pattern = re.compile(rf'\b{c_esc}\b', re.IGNORECASE)
    if not pattern.search(s):
        return None

    # Only blank out if the sentence is definitional or explanatory
    has_structure = re.search(
        r'\b(?:is|are|was|were|refers?\s+to|means?|called|known\s+as|'
        r'allows?|enables?|provides?|used\s+(?:to|for|in)|'
        r'defined\s+as|responsible\s+for|consists?\s+of)\b',
        s, re.IGNORECASE
    )
    if not has_structure:
        return None

    blanked = pattern.sub('________', s, count=1)
    if blanked == s:
        return None

    question = f"Fill in the blank: {blanked}"
    return {'question': question, 'answer': concept}


def _make_true_false_card(sentence: str) -> dict | None:
    """Create a True/False card — only from clear factual statements."""
    if _is_junk_sentence(sentence):
        return None
    s = sentence.strip()
    if not re.match(r'^[A-Z]', s):
        return None
    if '?' in s:
        return None
    if len(s) < 50 or len(s) > 250:
        return None
    # Must contain a clear factual verb structure
    if not re.search(r'\b(?:is|are|was|were|can|has|have|will|provides?|allows?|enables?)\b', s, re.I):
        return None

    question = f"True or False: {s}"
    answer = "True"
    return {'question': question, 'answer': answer}


def _attach_syntax(answer: str, concept: str, code_blocks: list[str]) -> str:
    """Append a relevant code snippet to the answer if one exists."""
    snippet = _find_syntax_for_concept(concept, code_blocks)
    if snippet:
        return f"{answer}\n\nSyntax:\n{snippet}"
    return answer


def generate_flashcards(file_path: str, file_type: str, max_cards: int = 20) -> list[dict]:
    """
    Main entry point: extract text from a file and generate flashcard Q&A pairs.
    Returns a list of {'question': ..., 'answer': ...} dicts.
    """
    raw_text = extract_text(file_path, file_type)
    if not raw_text or len(raw_text) < 50:
        return []

    # Extract code blocks from raw text BEFORE cleaning (cleaning collapses whitespace)
    code_blocks = _extract_code_blocks(raw_text)

    text = clean_text(raw_text)
    sentences = split_sentences(text)
    if not sentences:
        return []

    concepts = _extract_concepts(text, top_n=50)
    cards: list[dict] = []
    used_questions: set[str] = set()
    used_sentences: set[str] = set()

    # Phase 1: Definition cards — highest quality
    for concept in concepts:
        if len(cards) >= max_cards:
            break
        for sentence in sentences:
            if sentence in used_sentences:
                continue
            definition = _extract_definition(sentence, concept)
            if definition:
                question = _make_concept_question(concept)
                if question not in used_questions:
                    answer = _attach_syntax(definition, concept, code_blocks)
                    cards.append({'question': question, 'answer': answer})
                    used_questions.add(question)
                    used_sentences.add(sentence)
                    break

    # Phase 2: Explanation cards — "What does X do?"
    for concept in concepts:
        if len(cards) >= max_cards:
            break
        for sentence in sentences:
            if sentence in used_sentences:
                continue
            card = _make_explanation_card(sentence, concept)
            if card and card['question'] not in used_questions:
                card['answer'] = _attach_syntax(card['answer'], concept, code_blocks)
                cards.append(card)
                used_questions.add(card['question'])
                used_sentences.add(sentence)
                break

    # Phase 3: Fill-in-the-blank cards
    for concept in concepts:
        if len(cards) >= max_cards:
            break
        for sentence in sentences:
            if sentence in used_sentences:
                continue
            card = _make_fill_blank_card(sentence, concept)
            if card and card['question'] not in used_questions:
                card['answer'] = _attach_syntax(card['answer'], concept, code_blocks)
                cards.append(card)
                used_questions.add(card['question'])
                used_sentences.add(sentence)
                break

    # Phase 4: True/False cards from remaining good sentences
    for sentence in sentences:
        if len(cards) >= max_cards:
            break
        if sentence in used_sentences:
            continue
        card = _make_true_false_card(sentence)
        if card and card['question'] not in used_questions:
            cards.append(card)
            used_questions.add(card['question'])
            used_sentences.add(sentence)

    return cards[:max_cards]
